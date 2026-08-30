"""
Extract ALL text info from slides and save to JSON
Helps identify heading vs body text later
"""

import json
import os
from pathlib import Path
from pptx import Presentation
from app.utils.xml_helpers import q as _q, off_ext as _off_ext, local_name as _local


def rgb_to_hex(r, g, b):
    """Convert RGB to hex color"""
    return "#{:02x}{:02x}{:02x}".format(r, g, b)


def extract_color_from_run(run_element):
    """Extract hex color from text run"""
    rPr = run_element.find(_q("a:rPr"))
    if rPr is None:
        return None
    
    # Look for solid fill
    solidFill = rPr.find(_q("a:solidFill"))
    if solidFill is None:
        return None
    
    # Look for srgbClr (RGB color)
    srgb = solidFill.find(_q("a:srgbClr"))
    if srgb is not None:
        hex_val = srgb.get("val")
        if hex_val:
            return "#" + hex_val
    
    return None


def get_text_run_info(run_element, shape_position):
    """
    Extract complete info from ONE text run
    
    Returns: {
        "text": "My Text",
        "font_size_pt": 24,
        "bold": True,
        "font_name": "Arial",
        "color": "#3366FF",
        "italic": False,
        "underline": False,
        "position_top_left": (1000000, 500000),
        "position_center": (3000000, 1250000),
        "shape_position": (1000000, 500000),
        "shape_size": (4000000, 1500000)
    }
    """
    
    # Get text
    t_elem = run_element.find(_q("a:t"))
    if t_elem is None or not t_elem.text:
        return None
    
    text = t_elem.text.strip()
    if not text:
        return None
    
    # Get run properties
    rPr = run_element.find(_q("a:rPr"))
    
    info = {
        "text": text,
        "font_size_pt": None,
        "bold": False,
        "font_name": None,
        "color": None,
        "italic": False,
        "underline": False,
        "position_top_left": shape_position,
        "position_center": None
    }
    
    if rPr is not None:
        # Font size
        sz = rPr.get("sz")
        if sz:
            try:
                info["font_size_pt"] = int(int(sz) / 100)
            except:
                pass
        
        # Bold
        b = rPr.get("b")
        info["bold"] = (b == "1" or b == "true")
        
        # Italic
        i = rPr.get("i")
        info["italic"] = (i == "1" or i == "true")
        
        # Underline
        u = rPr.get("u")
        info["underline"] = (u and u != "none")
        
        # Font name
        latin = rPr.find(_q("a:latin"))
        if latin is not None:
            info["font_name"] = latin.get("typeface")
        
        # Color
        info["color"] = extract_color_from_run(run_element)
    
    # Calculate center position (approximate)
    if shape_position and shape_position[0] is not None:
        info["position_center"] = shape_position
    
    return info



HEADING_FONT_PT = 80
def pick_heading(text_elements, slide_height):
    """Runs that look like the slide title → one heading dict, or None."""
    if not text_elements:
        return None

    size_80 = [e for e in text_elements if e.get("font_size_pt") == HEADING_FONT_PT]
    if not size_80:
        return None

    def sort_key(e):
        pos = e.get("position_top_left") or (0, 10**12)
        y = pos[1] if len(pos) > 1 else 10**12
        x = pos[0] if pos else 0
        # top first, then left, bold preferred
        return (y, x, 0 if e.get("bold") else 1)

    pool = sorted(size_80, key=sort_key)
    first = pool[0]
    pos = first.get("position_top_left")
    same_shape = [
        e for e in pool
        if e.get("position_top_left") == pos
    ]
    # keep reading order within the same shape
    text = "".join(e["text"] for e in same_shape)

    return {
        "text": text,
        "font_size_pt": first.get("font_size_pt"),
        "bold": first.get("bold"),
        "font_name": first.get("font_name"),
        "color": first.get("color"),
        "italic": first.get("italic"),
        "underline": first.get("underline"),
        "position_top_left": first.get("position_top_left"),
        "position_center": first.get("position_center"),
        "shape_size": first.get("shape_size"),
        "used_size_80": bool(size_80),
        "run_count": len(same_shape),
    }


def extract_text_from_slide(slide, slide_number,prs):
    """
    Extract ALL text from one slide
    
    Returns: {
        "slide_number": 0,
        "slide_width": 9144000,
        "slide_height": 6858000,
        "text_elements": [
            {text info...},
            {text info...}
        ]
    }
    """
    
    sptree = slide._element.find(_q("p:cSld") + "/" + _q("p:spTree"))
    if sptree is None:
        return {
            "slide_number": slide_number,
            "text_elements": []
        }
    
    # Get all shapes
    children = [c for c in list(sptree) if _local(c) == "sp"]
    
    text_elements = []
    
    for child in children:
        # Get shape position and size
        off, ext = _off_ext(child, "p:spPr")
        
        # Look for text body
        txBody = child.find(_q("p:txBody"))
        if txBody is None:
            continue
        
        # Loop through all text runs
        for para in txBody.iter(_q("a:p")):
            for run in para.iter(_q("a:r")):
                text_info = get_text_run_info(run, off)
                
                if text_info:
                    # Add shape size info
                    text_info["shape_size"] = ext
                    text_elements.append(text_info)
    if sptree is None:
        return {
            "slide_number": slide_number,
            "heading": None,
        }

    # ... same loop into text_elements ...
    heading = pick_heading(text_elements, int(prs.slide_height))

    return {
        "slide_number": slide_number,
        "slide_width": int(prs.slide_width),
        "slide_height": int(prs.slide_height),
        "heading": heading,
    }


def analyze_slides_text(pptx_path, start_index=None, end_index=None):
    """
    Extract text info from slides and save to JSON files
    
    Args:
        pptx_path: Path to PPTX file (e.g., "input.pptx")
        start_index: Start slide number (0-based). Default: 0
        end_index: End slide number (inclusive). Default: None (all slides)
    
    Examples:
        # Extract ALL slides
        analyze_slides_text("input.pptx")
        
        # Extract slides 0-5
        analyze_slides_text("input.pptx", start_index=0, end_index=5)
        
        # Extract only slide 3
        analyze_slides_text("input.pptx", start_index=3, end_index=3)
        
        # Extract from slide 2 onwards
        analyze_slides_text("input.pptx", start_index=2)
    """
    
    # Load presentation
    prs = Presentation(pptx_path)
    
    # Set default indices
    if start_index is None:
        start_index = 0
    if end_index is None:
        end_index = len(prs.slides) - 1
    
    # Ensure valid range
    start_index = max(0, start_index)
    end_index = min(len(prs.slides) - 1, end_index)
    
    print(f"Analyzing slides {start_index} to {end_index}...")
    
    # Create output directory
    output_dir = Path("app/data/slide_text_analysis")
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # Process each slide
    for idx, slide in enumerate(prs.slides):
        if idx < start_index or idx > end_index:
            continue
        
        print(f"  Processing slide {idx}...")
        
        # Extract text
        slide_data = extract_text_from_slide(slide, idx,prs)
        
        # Save to JSON
        output_file = output_dir / f"slide_{idx}.json"
        with open(output_file, "w", encoding="utf-8") as f:
            json.dump(slide_data, f, indent=2, ensure_ascii=False)
        
        print(f"    Saved: {output_file}")
    
    print(f"Done! Files saved to: {output_dir}")
    return output_dir