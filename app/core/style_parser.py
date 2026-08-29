from __future__ import annotations

import zipfile
from dataclasses import dataclass, field
from io import BytesIO
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

from app.constants.xml_namespaces import A, NS, P


@dataclass
class ShapeStyle:
    name: str
    shape_type: str
    left: int | None
    top: int | None
    width: int | None
    height: int | None
    has_text: bool
    sample_text: str
    fill_hex: str | None
    font_name: str | None
    font_size_pt: int | None
    font_color_hex: str | None
    is_group: bool = False
    is_table: bool = False
    is_picture: bool = False

    @property
    def role(self) -> str:
        """Heuristic role label used by AI matching."""
        n = self.name.lower() + " " + self.sample_text.lower()
        if "question" in n:
            return "question_badge"
        if self.is_picture:
            return "picture"
        if self.is_table:
            return "table"
        if self.is_group and self.width and self.height and self.width < 2000000:
            return "option_pill"
        if self.is_group:
            return "decorative_group"
        if self.has_text and self.width and self.width > 30000000 and self.height and self.height < 1500000:
            return "heading"
        if self.has_text and self.sample_text:
            return "body"
        return "decoration"


@dataclass
class SlideDesign:
    index: int
    layout_name: str
    background_xml: bytes | None
    background_kind: str  # "image" | "solid" | "none"
    theme_font: str | None
    theme_colors: dict[str, str] = field(default_factory=dict)
    shapes: list[ShapeStyle] = field(default_factory=list)

    @property
    def fingerprint(self) -> dict[str, Any]:
        return {
            "layout": self.layout_name,
            "shape_count": len(self.shapes),
            "text_count": sum(1 for s in self.shapes if s.has_text),
            "group_count": sum(1 for s in self.shapes if s.is_group),
            "table_count": sum(1 for s in self.shapes if s.is_table),
            "picture_count": sum(1 for s in self.shapes if s.is_picture),
            "option_pills": sum(1 for s in self.shapes if s.role == "option_pill"),
            "has_question_badge": any(s.role == "question_badge" for s in self.shapes),
            "has_heading": any(s.role == "heading" for s in self.shapes),
            "has_body": any(s.role == "body" for s in self.shapes),
            "has_decorative_group": any(s.role == "decorative_group" for s in self.shapes),
        }


def _xml_bytes(elem: etree._Element | None) -> bytes | None:
    if elem is None:
        return None
    return etree.tostring(elem)


def _bg_kind(bg_elem: etree._Element | None) -> str:
    if bg_elem is None:
        return "none"
    blip = bg_elem.find(".//{%s}blip" % A)
    if blip is not None:
        return "image"
    solid = bg_elem.find(".//{%s}solidFill" % A)
    if solid is not None:
        return "solid"
    return "other"


def _parse_hex(elem: etree._Element | None) -> str | None:
    if elem is None:
        return None
    srgb = elem.find("{%s}srgbClr" % A)
    if srgb is not None:
        return srgb.get("val", "").upper()
    scheme = elem.find("{%s}schemeClr" % A)
    if scheme is not None:
        return f"scheme:{scheme.get('val', '')}"
    return None


def _shape_style(sp_elem: etree._Element, spPr: etree._Element | None) -> dict[str, Any]:
    """Extract a flat dict of style fields from a shape XML element."""
    info: dict[str, Any] = {
        "name": sp_elem.get("name", ""),
        "is_group": False,
        "is_table": False,
        "is_picture": False,
        "has_text": False,
        "sample_text": "",
        "fill_hex": None,
        "font_name": None,
        "font_size_pt": None,
        "font_color_hex": None,
    }

    tag = etree.QName(sp_elem).localname
    if tag == "grpSp":
        info["is_group"] = True
    elif tag == "graphicFrame":
        # table or chart container
        graphic = sp_elem.find(".//{%s}graphicData" % A)
        uri = graphic.get("uri", "") if graphic is not None else ""
        if "table" in uri.lower():
            info["is_table"] = True
    elif tag == "pic":
        info["is_picture"] = True

    if spPr is not None:
        solid = spPr.find("{%s}solidFill" % A)
        info["fill_hex"] = _parse_hex(solid)

    txBody = sp_elem.find("{%s}txBody" % P)
    if txBody is not None:
        info["has_text"] = True
        info["sample_text"] = "".join(t.text or "" for t in txBody.iter("{%s}t" % A)).strip()[:80]
        first_r = txBody.find(".//{%s}r" % A)
        if first_r is not None:
            rPr = first_r.find("{%s}rPr" % A)
            if rPr is not None:
                sz = rPr.get("sz")
                if sz:
                    try:
                        info["font_size_pt"] = int(int(sz) / 100)
                    except ValueError:
                        pass
                latin = rPr.find("{%s}latin" % A)
                if latin is not None:
                    info["font_name"] = latin.get("typeface")
                solid = rPr.find("{%s}solidFill" % A)
                info["font_color_hex"] = _parse_hex(solid)
            else:
                defRPr = txBody.find("{%s}lstStyle/{%s}defRPr" % (P, A))
                if defRPr is not None:
                    sz = defRPr.get("sz")
                    if sz:
                        try:
                            info["font_size_pt"] = int(int(sz) / 100)
                        except ValueError:
                            pass
                    latin = defRPr.find("{%s}latin" % A)
                    if latin is not None:
                        info["font_name"] = latin.get("typeface")

    return info


def _geom(spPr: etree._Element | None) -> tuple[int | None, int | None, int | None, int | None]:
    if spPr is None:
        return None, None, None, None
    xfrm = spPr.find("{%s}xfrm" % A) or spPr.find("{%s}off" % A)
    off = spPr.find("{%s}xfrm/{%s}off" % (A, A))
    ext = spPr.find("{%s}xfrm/{%s}ext" % (A, A))
    if off is None or ext is None:
        return None, None, None, None
    return (
        int(off.get("x", 0)),
        int(off.get("y", 0)),
        int(ext.get("cx", 0)),
        int(ext.get("cy", 0)),
    )


def _extract_theme(zip_f: zipfile.ZipFile) -> tuple[str | None, dict[str, str]]:
    """Return (latin_major, {scheme_role: hex}) from the first theme part found."""
    theme_names = [n for n in zip_f.namelist() if n.startswith("ppt/theme/theme") and n.endswith(".xml")]
    if not theme_names:
        return None, {}
    xml = zip_f.read(theme_names[0])
    root = etree.fromstring(xml)
    scheme = root.find(".//{%s}clrScheme" % A)
    colors: dict[str, str] = {}
    if scheme is not None:
        for child in scheme:
            role = etree.QName(child).localname
            srgb = child.find("{%s}srgbClr" % A)
            if srgb is not None:
                colors[role] = srgb.get("val", "").upper()
            else:
                sys_clr = child.find("{%s}sysClr" % A)
                if sys_clr is not None:
                    colors[role] = "sys:" + (sys_clr.get("lastClr") or sys_clr.get("val") or "")

    major = root.find(".//{%s}majorFont/{%s}latin" % (A, A))
    minor = root.find(".//{%s}minorFont/{%s}latin" % (A, A))
    font = None
    if major is not None and major.get("typeface"):
        font = major.get("typeface")
    elif minor is not None and minor.get("typeface"):
        font = minor.get("typeface")
    return font, colors


def parse_template(template_path: str) -> tuple[Presentation, list[SlideDesign], dict[str, bytes]]:
    """Parse the template. Return (prs, designs, archive_bytes_for_clone)."""
    with open(template_path, "rb") as fh:
        archive_bytes = fh.read()
    zip_f = zipfile.ZipFile(BytesIO(archive_bytes))
    theme_font, theme_colors = _extract_theme(zip_f)
    zip_f.close()

    prs = Presentation(template_path)
    designs: list[SlideDesign] = []
    for idx, slide in enumerate(prs.slides):
        sxml = slide._element
        cSld = sxml.find("{%s}cSld" % P)
        bg_elem = cSld.find("{%s}bg" % P) if cSld is not None else None

        shapes: list[ShapeStyle] = []
        for sp in sxml.iter():
            if etree.QName(sp).localname not in {"sp", "pic", "grpSp", "graphicFrame"}:
                continue
            spPr = sp.find("{%s}spPr" % P)
            left, top, w, h = _geom(spPr)
            info = _shape_style(sp, spPr)
            shapes.append(
                ShapeStyle(
                    name=info["name"],
                    shape_type=etree.QName(sp).localname,
                    left=left,
                    top=top,
                    width=w,
                    height=h,
                    has_text=info["has_text"],
                    sample_text=info["sample_text"],
                    fill_hex=info["fill_hex"],
                    font_name=info["font_name"],
                    font_size_pt=info["font_size_pt"],
                    font_color_hex=info["font_color_hex"],
                    is_group=info["is_group"],
                    is_table=info["is_table"],
                    is_picture=info["is_picture"],
                )
            )

        designs.append(
            SlideDesign(
                index=idx,
                layout_name=slide.slide_layout.name,
                background_xml=_xml_bytes(bg_elem),
                background_kind=_bg_kind(bg_elem),
                theme_font=theme_font,
                theme_colors=theme_colors,
                shapes=shapes,
            )
        )

    return prs, designs, archive_bytes


def parse_input_slide_signature(input_path: str) -> list[dict[str, Any]]:
    """Lightweight structural signature per input slide, used to pre-match
    slides before asking Azure AI for confirmation."""
    prs = Presentation(input_path)
    sigs: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides):
        text_count = 0
        group_count = 0
        table_count = 0
        picture_count = 0
        option_pill_count = 0
        question_badge = False
        body_text_count = 0
        sample = ""
        for sp in slide.shapes:
            tag = sp.shape_type
            if tag is None:
                group_count += 1
                try:
                    cw = sp.width or 0
                    ch = sp.height or 0
                    if cw and ch and cw < 2000000 and ch < 2000000:
                        option_pill_count += 1
                except Exception:
                    pass
                continue
            if hasattr(tag, "name") and tag.name == "TABLE":
                table_count += 1
            elif hasattr(tag, "name") and tag.name == "PICTURE":
                picture_count += 1
            elif sp.has_text_frame:
                text_count += 1
                txt = sp.text_frame.text.strip().lower()
                if not sample:
                    sample = sp.text_frame.text.strip()[:120]
                if "question" in txt or "questions" in txt:
                    question_badge = True
                if len(txt) > 30:
                    body_text_count += 1
        sigs.append(
            {
                "index": idx,
                "shape_count": len(slide.shapes),
                "text_count": text_count,
                "group_count": group_count,
                "table_count": table_count,
                "picture_count": picture_count,
                "option_pills": option_pill_count,
                "has_question_badge": question_badge,
                "body_text_count": body_text_count,
                "sample": sample,
            }
        )
    return sigs


def emu_to_in(emu: int | None) -> float | None:
    return None if emu is None else Emu(emu).inches


def hex_no_hash(value: str | None) -> str | None:
    if value is None:
        return None
    if value.startswith("scheme:"):
        return None
    return value.lstrip("#").upper()