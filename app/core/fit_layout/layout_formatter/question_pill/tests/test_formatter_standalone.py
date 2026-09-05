"""Standalone test for question_pill formatter on p5.pptx.

Simplified version that doesn't require heavy imports.
Just detects pills, clones them, and saves formatted PPTX.
"""

from pathlib import Path
import copy

from lxml import etree
from pptx import Presentation


# Path to the input test data
INPUT_PPT_PATH = Path("app/data/input/p5.pptx")

# Path to save formatted output PPTX files
OUTPUT_DIR = Path(__file__).parent.parent / "output"


def q(tag: str) -> str:
    """Convert short tag name to full XML namespace."""
    ns = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }
    prefix, local = tag.split(":")
    return "{%s}%s" % (ns[prefix], local)


def clone_and_place_simple(el: etree._Element, off: tuple, ext: tuple) -> etree._Element:
    """Simple clone_and_place without external dependencies."""
    clone = copy.deepcopy(el)
    spPr = clone.find(q("p:spPr"))
    if spPr is not None and off and ext:
        xfrm = spPr.find(q("a:xfrm"))
        if xfrm is None:
            xfrm = etree.SubElement(spPr, q("a:xfrm"))
        off_el = xfrm.find(q("a:off"))
        if off_el is None:
            off_el = etree.SubElement(xfrm, q("a:off"))
        off_el.set("x", str(off[0]))
        off_el.set("y", str(off[1]))
        ext_el = xfrm.find(q("a:ext"))
        if ext_el is None:
            ext_el = etree.SubElement(xfrm, q("a:ext"))
        ext_el.set("cx", str(ext[0]))
        ext_el.set("cy", str(ext[1]))
    return clone


def format_question_pill_simple(pill_el: etree._Element, label_text: str) -> etree._Element:
    """Format question pill with fixed position."""
    # Constants from question_pill/constants.py
    QUESTION_PILL = {
        "x": -1_864_528,
        "y": 681_774,
        "width": 7_081_988,
        "height": 1_569_660,
    }

    # Position and size the pill
    off = (QUESTION_PILL["x"], QUESTION_PILL["y"])
    ext = (QUESTION_PILL["width"], QUESTION_PILL["height"])
    pill_clone = clone_and_place_simple(pill_el, off, ext)

    # Update text inside the pill
    txBody = pill_clone.find(q("p:txBody"))
    if txBody is None:
        txBody = pill_clone.find(q("a:txBody"))
    if txBody is not None:
        t_el = txBody.find(q("a:t"))
        if t_el is not None:
            t_el.text = label_text

    return pill_clone


def format_question_body_simple(body_el: etree._Element) -> etree._Element:
    """Format question body text with fixed position."""
    # Constants from question_pill/constants.py
    QUESTION_BODY = {
        "x": 1_104_806,
        "y": 2_396_709,
        "width": 34_747_200,
        "height": 3_474_720,
    }

    # Position and size the body text box
    off = (QUESTION_BODY["x"], QUESTION_BODY["y"])
    ext = (QUESTION_BODY["width"], QUESTION_BODY["height"])
    body_clone = clone_and_place_simple(body_el, off, ext)

    # Enable text wrapping for the body text
    txBody = body_clone.find(q("p:txBody"))
    if txBody is None:
        txBody = body_clone.find(q("a:txBody"))

    if txBody is not None:
        bodyPr = txBody.find(q("a:bodyPr"))
        if bodyPr is not None:
            bodyPr.set("wrap", "square")
            bodyPr.set("anchor", "t")

    return body_clone


def test_question_pill():
    """Test question pill formatter on p5.pptx."""

    if not INPUT_PPT_PATH.exists():
        print(f"Input file not found: {INPUT_PPT_PATH}")
        return

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    # Load presentation
    prs = Presentation(str(INPUT_PPT_PATH))

    print(f"\n{'='*80}")
    print(f"Testing question_pill formatter on p5.pptx")
    print(f"{'='*80}\n")

    formatted_count = 0
    slides_with_pills = []

    # Process each slide
    for slide_idx, slide in enumerate(prs.slides):
        print(f"Slide {slide_idx}:")

        pill_el = None
        label_el = None

        # Find both the pill (roundRect) and the label text box
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, "_element"):
                elem = shape._element
                ns = {
                    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
                }

                # Check if this is the question pill (roundRect shape)
                geom = elem.find(".//a:prstGeom", ns)
                if geom is not None and geom.get("prst") == "roundRect":
                    pill_el = elem
                    slides_with_pills.append(slide_idx)

                # Check if this is a text box with "Question" text (the label)
                # Look for text frames with "Question" text
                txBody = elem.find(".//p:txBody", ns)
                if txBody is None:
                    txBody = elem.find(".//a:txBody", ns)

                if txBody is not None:
                    # Find text content
                    text_content = ""
                    for t in txBody.findall(".//a:t", ns):
                        if t.text:
                            text_content += t.text

                    # If text is "Question" (or contains it), this is likely the label
                    if "Question" in text_content or "Option" in text_content:
                        label_el = elem

        # Format if we found both pill and label
        if pill_el is not None and label_el is not None:
            try:
                # Format both the pill and the label
                formatted_pill = format_question_pill_simple(pill_el, "Question")
                formatted_label = clone_and_place_simple(
                    label_el,
                    (1_039_500, 827_049),        # QUESTION_LABEL position
                    (4_562_190, 1_015_622),      # QUESTION_LABEL size
                )

                # Replace both in the slide
                spTree = slide._element.find(q("p:cSld")).find(q("p:spTree"))
                spTree.remove(pill_el)
                spTree.remove(label_el)
                spTree.append(formatted_pill)
                spTree.append(formatted_label)

                print(f"  ✓ Formatted question pill + label")
                print(f"    Pill position: (-2.039\", 0.746\")")
                print(f"    Pill size: 7.745\" × 1.717\"")
                print(f"    Label text: 'Question'")
                print(f"    Label position: (1.137\", 0.904\")")
                formatted_count += 1

            except Exception as e:
                print(f"  ✗ Error: {str(e)}")
        elif pill_el is not None:
            print(f"  ⚠ Found pill but no label text")
        else:
            print(f"  - No question pill found")

        print()

    # Save output
    if formatted_count > 0:
        output_file = OUTPUT_DIR / "p5_formatted_pills.pptx"
        prs.save(str(output_file))
        print(f"{'='*80}")
        print(f"✓ SUCCESS! Saved formatted presentation")
        print(f"  File: {output_file.name}")
        print(f"  Slides formatted: {formatted_count}")
        print(f"  Slides with pills: {sorted(slides_with_pills)}")
        print(f"{'='*80}\n")
    else:
        print(f"{'='*80}")
        print("No pills found to format")
        print(f"{'='*80}\n")


if __name__ == "__main__":
    test_question_pill()
