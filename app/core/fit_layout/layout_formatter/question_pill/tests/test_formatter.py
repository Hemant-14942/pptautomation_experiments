"""Test the question_pill formatter on p5.pptx input slides.

Tests the format_question_pill() function by:
1. Loading p5.pptx (which contains question pill slides)
2. Finding slides with question pill shapes
3. Cloning each pill and positioning it with the formatter
4. Saving formatted slides to PPTX files in the output folder
"""

from pathlib import Path

from lxml import etree
from pptx import Presentation

from app.core.fit_layout.layout_formatter.question_pill.formatter import (
    format_question_pill,
)
from app.utils.xml_helpers import q


# Path to the input test data
INPUT_PPT_PATH = Path("app/data/input/p5.pptx")

# Path to save formatted output PPTX files for manual inspection
# Output goes to question_pill/output/, not inside tests/
OUTPUT_DIR = Path(__file__).parent.parent / "output"


def test_question_pill_formatter():
    """Test formatting question pills from p5.pptx slides."""

    # Check if input file exists
    if not INPUT_PPT_PATH.exists():
        print(f"Input file not found: {INPUT_PPT_PATH}")
        return

    # Create output directory if it doesn't exist
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    # Load the presentation
    prs = Presentation(str(INPUT_PPT_PATH))

    print(f"\n{'='*80}")
    print(f"Testing question_pill formatter on {INPUT_PPT_PATH}")
    print(f"{'='*80}\n")

    # Track results
    formatted_count = 0
    slides_with_pills = []

    # Check each slide for question pill shapes
    for slide_idx, slide in enumerate(prs.slides):
        print(f"Slide {slide_idx}:")

        # Look for AUTO_SHAPE elements (which include the question pill)
        question_pill_shape = None
        question_pill_shape_idx = None

        for shape_idx, shape in enumerate(slide.shapes):
            # Check if this is a rounded rectangle (the question pill)
            if hasattr(shape, "_element"):
                elem = shape._element
                ns = {
                    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
                }
                geom = elem.find(".//a:prstGeom", ns)
                if geom is not None and geom.get("prst") == "roundRect":
                    question_pill_shape = shape
                    question_pill_shape_idx = shape_idx
                    slides_with_pills.append(slide_idx)
                    break

        if question_pill_shape is not None:
            # Format the pill using our formatter
            try:
                # Get the pill element from the slide
                pill_el = question_pill_shape._element

                # Format the pill with label text
                label_text = "Question"  # Default label
                formatted_pill = format_question_pill(pill_el, label_text)

                # Replace the original pill with the formatted one in the slide
                # Find the parent shape tree (spTree) that contains all shapes
                spTree = slide._element.find(q("p:cSld")).find(q("p:spTree"))

                # Remove the original pill shape from the slide
                spTree.remove(pill_el)

                # Add the formatted pill to the slide
                spTree.append(formatted_pill)

                print(f"  ✓ Formatted question pill at shape {question_pill_shape_idx}")
                print(f"    Position: ({question_pill_shape.left}, {question_pill_shape.top}) EMU")
                print(f"    Label set to: '{label_text}'")

                formatted_count += 1

            except Exception as e:
                print(f"  ✗ Error formatting pill: {str(e)}")

        else:
            print(f"  - No question pill found")

        print()

    # Save the modified presentation to output
    if formatted_count > 0:
        output_pptx = OUTPUT_DIR / f"p5_formatted_pills.pptx"
        prs.save(str(output_pptx))
        print(f"{'='*80}")
        print(f"✓ Saved formatted presentation: {output_pptx.name}")
        print(f"  Formatted slides: {formatted_count}")
        print(f"  Slides with pills: {sorted(slides_with_pills)}")
        print(f"{'='*80}\n")
    else:
        print(f"{'='*80}")
        print("No slides with question pills found to format.")
        print(f"{'='*80}\n")


if __name__ == "__main__":
    test_question_pill_formatter()
