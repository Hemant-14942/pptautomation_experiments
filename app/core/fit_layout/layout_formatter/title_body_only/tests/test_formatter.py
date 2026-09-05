"""Diagnostic check: does a given (pptx_path, slide_index) qualify as
title_body_only, and if so, does the formatter place its body text correctly?

Run directly for a manual check on one slide:
    python -m app.core.fit_layout.layout_formatter.title_body_only.tests.test_formatter <pptx_path> <slide_index>
"""

import os
import sys

from pptx import Presentation
from pptx.util import Emu

from app.core.slide_text_analyzer import extract_text_from_slide
from app.core.style_parser import parse_input_slide_signature
from app.core.fit_layout.layout_formatter.slide_type_detector import (
    SLIDE_TYPE_TITLE_BODY_ONLY,
    detect_slide_type,
)
from app.core.fit_layout.layout_formatter.title_body_only.constants import BODY_TEXTBOX
from app.core.fit_layout.layout_formatter.title_body_only.formatter import format_body_only

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "output")

XFRM_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def _get_off_ext(el):
    xfrm = el.find(".//p:spPr/a:xfrm", XFRM_NS)
    off = xfrm.find("a:off", XFRM_NS)
    ext = xfrm.find("a:ext", XFRM_NS)
    return (int(off.get("x")), int(off.get("y"))), (int(ext.get("cx")), int(ext.get("cy")))


def _find_body_el(slide, prs):
    """Body text shape (>30 chars), excluding the heading's own shape so a
    long heading is never mistaken for body text."""
    heading = extract_text_from_slide(slide, 0, prs).get("heading")
    heading_pos = heading.get("position_top_left") if heading else None
    for sp in slide.shapes:
        if not sp.has_text_frame:
            continue
        if heading_pos is not None and (sp.left, sp.top) == heading_pos:
            continue
        if len(sp.text_frame.text.strip()) > 30:
            return sp._element
    return None


def _get_anchor(el):
    xfrm_ns = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }
    txBody = el.find(".//p:txBody", xfrm_ns)
    if txBody is None:
        return None
    bodyPr = txBody.find("a:bodyPr", xfrm_ns)
    if bodyPr is None:
        return None
    return bodyPr.get("anchor")


def verify_slide(pptx_path: str, slide_index: int) -> dict:
    """Check one slide end-to-end: detection + formatter placement.

    Returns a report dict — "pass": bool plus details for inspection.
    Does not raise on a non-matching slide; it reports why.
    """
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_index]
    sig = parse_input_slide_signature(pptx_path)[slide_index]

    detected_type = detect_slide_type(sig, slide, prs)
    if detected_type != SLIDE_TYPE_TITLE_BODY_ONLY:
        return {
            "pass": False,
            "reason": f"slide type detected as {detected_type!r}, not title_body_only",
            "detected_type": detected_type,
        }

    body_el = _find_body_el(slide, prs)

    if body_el is None:
        return {
            "pass": False,
            "reason": "detector matched but could not locate body text shape",
            "detected_type": detected_type,
        }

    body_clone = format_body_only(body_el)
    body_off, body_ext = _get_off_ext(body_clone)
    anchor = _get_anchor(body_clone)

    body_ok = (body_off, body_ext) == (
        (BODY_TEXTBOX["x"], BODY_TEXTBOX["y"]),
        (BODY_TEXTBOX["width"], BODY_TEXTBOX["height"]),
    )
    anchor_ok = anchor == "ctr"

    return {
        "pass": body_ok and anchor_ok,
        "detected_type": detected_type,
        "body": {
            "off": body_off,
            "ext": body_ext,
            "matches_fixed_box": body_ok,
        },
        "anchor": {"value": anchor, "is_center": anchor_ok},
    }


def save_formatted_slide(pptx_path: str, slide_index: int, output_dir: str = OUTPUT_DIR) -> str | None:
    """Apply the formatter to one slide in-place and save the whole
    presentation as a new pptx in output_dir. Returns the saved path, or
    None if the slide doesn't match title_body_only."""
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_index]
    sig = parse_input_slide_signature(pptx_path)[slide_index]

    if detect_slide_type(sig, slide, prs) != SLIDE_TYPE_TITLE_BODY_ONLY:
        return None

    body_el = _find_body_el(slide, prs)

    if body_el is None:
        return None

    body_clone = format_body_only(body_el)
    spTree = slide.shapes._spTree
    spTree.remove(body_el)
    spTree.append(body_clone)

    os.makedirs(output_dir, exist_ok=True)
    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    out_path = os.path.join(output_dir, f"{base_name}_slide{slide_index + 1}_formatted.pptx")
    prs.save(out_path)
    return out_path


def _print_report(pptx_path: str, slide_index: int, report: dict) -> None:
    print(f"\n{pptx_path} | slide index {slide_index} (PPT slide {slide_index + 1})")
    print(f"  detected_type: {report.get('detected_type')}")
    if not report["pass"] and "reason" in report:
        print(f"  RESULT: FAIL — {report['reason']}")
        return

    b = report["body"]
    a = report["anchor"]
    b_off_in = (Emu(b["off"][0]).inches, Emu(b["off"][1]).inches)
    b_ext_in = (Emu(b["ext"][0]).inches, Emu(b["ext"][1]).inches)

    print(f"  body:   off={b_off_in[0]:.2f},{b_off_in[1]:.2f}in  ext={b_ext_in[0]:.2f}x{b_ext_in[1]:.2f}in  fixed_box_match={b['matches_fixed_box']}")
    print(f"  anchor: {a['value']!r}  center={a['is_center']}")
    print(f"  RESULT: {'PASS' if report['pass'] else 'FAIL'}")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python -m ...test_formatter <pptx_path> <slide_index>")
        sys.exit(1)
    path, idx = sys.argv[1], int(sys.argv[2])
    result = verify_slide(path, idx)
    _print_report(path, idx, result)
    if result["pass"]:
        saved_path = save_formatted_slide(path, idx)
        print(f"  saved formatted pptx: {saved_path}")
