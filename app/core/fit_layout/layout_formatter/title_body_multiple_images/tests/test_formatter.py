"""Diagnostic check for title_body_multiple_images slides.

Run directly to test one slide:
    python -m app.core.fit_layout.layout_formatter.title_body_multiple_images.tests.test_formatter <pptx_path> <slide_index>
"""

import os
import sys

from pptx import Presentation
from pptx.util import Emu

from app.core.slide_text_analyzer import extract_text_from_slide
from app.core.style_parser import parse_input_slide_signature
from app.core.fit_layout.layout_formatter.slide_type_detector import (
    IMAGE_Y_THRESHOLD_EMU,
    SLIDE_TYPE_TITLE_BODY_MULTIPLE_IMAGES,
    detect_slide_type,
)
from app.core.fit_layout.layout_formatter.title_body_multiple_images.constants import (
    BODY_TEXTBOX,
    get_image_boxes,
)
from app.core.fit_layout.layout_formatter.title_body_multiple_images.formatter import (
    format_title_body_multiple_images,
)

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "output")

XFRM_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def _get_off_ext(el):
    xfrm = el.find(".//p:spPr/a:xfrm", XFRM_NS)
    off = xfrm.find("a:off", XFRM_NS)
    ext = xfrm.find("a:ext", XFRM_NS)
    return (int(off.get("x")), int(off.get("y"))), (int(ext.get("cx")), int(ext.get("cy")))


def _find_body_and_pictures(slide, prs):
    """Body text shape (>30 chars, excluding the heading's own shape) and
    all real content pictures (below the header area)."""
    heading = extract_text_from_slide(slide, 0, prs).get("heading")
    heading_pos = heading.get("position_top_left") if heading else None
    body_el = None
    pic_els = []
    pic_sizes = []
    for sp in slide.shapes:
        stype = sp.shape_type
        if stype is not None and getattr(stype, "name", "") == "PICTURE" and sp.top and sp.top > IMAGE_Y_THRESHOLD_EMU:
            pic_els.append(sp._element)
            pic_sizes.append(sp.image.size)
        elif sp.has_text_frame and (sp.left, sp.top) != heading_pos and len(sp.text_frame.text.strip()) > 30:
            body_el = sp._element
    return body_el, pic_els, pic_sizes


def verify_slide(pptx_path: str, slide_index: int) -> dict:
    """Check if slide qualifies and test formatter.

    Returns a report dict with detection + formatting results.
    """
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_index]
    sig = parse_input_slide_signature(pptx_path)[slide_index]

    detected_type = detect_slide_type(sig, slide, prs)
    if detected_type != SLIDE_TYPE_TITLE_BODY_MULTIPLE_IMAGES:
        return {
            "pass": False,
            "reason": f"detected as {detected_type!r}, not title_body_multiple_images",
            "detected_type": detected_type,
        }

    body_el, pic_els, pic_sizes = _find_body_and_pictures(slide, prs)

    if body_el is None or not pic_els:
        return {
            "pass": False,
            "reason": "detector matched but could not re-locate body/picture shapes",
            "detected_type": detected_type,
        }

    body_clone, pic_clones = format_title_body_multiple_images(body_el, pic_els, pic_sizes)

    body_off, body_ext = _get_off_ext(body_clone)
    body_ok = (body_off, body_ext) == (
        (BODY_TEXTBOX["x"], BODY_TEXTBOX["y"]),
        (BODY_TEXTBOX["width"], BODY_TEXTBOX["height"]),
    )

    expected_boxes = get_image_boxes(len(pic_els))
    pic_results = []
    all_pics_ok = True

    for pic_clone, pic_size, expected_box in zip(pic_clones, pic_sizes, expected_boxes):
        pic_off, pic_ext = _get_off_ext(pic_clone)

        expected_aspect = pic_size[0] / pic_size[1]
        actual_aspect = pic_ext[0] / pic_ext[1]
        aspect_ok = abs(actual_aspect - expected_aspect) < 0.01

        within_box = (
            pic_off[0] >= expected_box["x"]
            and pic_off[1] >= expected_box["y"]
            and pic_off[0] + pic_ext[0] <= expected_box["x"] + expected_box["width"]
            and pic_off[1] + pic_ext[1] <= expected_box["y"] + expected_box["height"]
        )

        pic_ok = aspect_ok and within_box
        all_pics_ok = all_pics_ok and pic_ok
        pic_results.append(
            {
                "off": pic_off,
                "ext": pic_ext,
                "aspect_ok": aspect_ok,
                "within_box": within_box,
                "pic_ok": pic_ok,
            }
        )

    return {
        "pass": body_ok and all_pics_ok,
        "detected_type": detected_type,
        "num_images": len(pic_els),
        "body": {"off": body_off, "ext": body_ext, "matches_fixed_box": body_ok},
        "pictures": pic_results,
    }


def save_formatted_slide(pptx_path: str, slide_index: int, output_dir: str = OUTPUT_DIR) -> str | None:
    """Apply formatter and save the full presentation."""
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_index]
    sig = parse_input_slide_signature(pptx_path)[slide_index]

    if detect_slide_type(sig, slide, prs) != SLIDE_TYPE_TITLE_BODY_MULTIPLE_IMAGES:
        return None

    body_el, pic_els, pic_sizes = _find_body_and_pictures(slide, prs)

    if body_el is None or not pic_els:
        return None

    body_clone, pic_clones = format_title_body_multiple_images(body_el, pic_els, pic_sizes)

    spTree = slide.shapes._spTree
    spTree.remove(body_el)
    for pic_el in pic_els:
        spTree.remove(pic_el)
    spTree.append(body_clone)
    for pic_clone in pic_clones:
        spTree.append(pic_clone)

    os.makedirs(output_dir, exist_ok=True)
    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    out_path = os.path.join(output_dir, f"{base_name}_slide{slide_index + 1}_formatted.pptx")
    prs.save(out_path)
    return out_path


def _print_report(pptx_path: str, slide_index: int, report: dict) -> None:
    print(f"\n{pptx_path} | slide index {slide_index} (PPT slide {slide_index + 1})")
    print(f"  detected_type: {report.get('detected_type')}")
    if not report["pass"] and "reason" in report:
        print(f"  RESULT: FAIL — {report['reason']}")
        return

    b = report["body"]
    b_off_in = (Emu(b["off"][0]).inches, Emu(b["off"][1]).inches)
    b_ext_in = (Emu(b["ext"][0]).inches, Emu(b["ext"][1]).inches)

    print(f"  body:  off={b_off_in[0]:.2f},{b_off_in[1]:.2f}in  ext={b_ext_in[0]:.2f}x{b_ext_in[1]:.2f}in  fixed_box={b['matches_fixed_box']}")
    print(f"  images: {report['num_images']} images")
    for i, p in enumerate(report["pictures"]):
        i_off_in = (Emu(p["off"][0]).inches, Emu(p["off"][1]).inches)
        i_ext_in = (Emu(p["ext"][0]).inches, Emu(p["ext"][1]).inches)
        print(f"    image {i+1}: off={i_off_in[0]:.2f},{i_off_in[1]:.2f}in  ext={i_ext_in[0]:.2f}x{i_ext_in[1]:.2f}in  aspect_ok={p['aspect_ok']}  within_box={p['within_box']}")

    print(f"  RESULT: {'PASS' if report['pass'] else 'FAIL'}")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python -m ...test_formatter <pptx_path> <slide_index>")
        sys.exit(1)
    path, idx = sys.argv[1], int(sys.argv[2])
    result = verify_slide(path, idx)
    _print_report(path, idx, result)
    if result["pass"]:
        saved_path = save_formatted_slide(path, idx)
        print(f"  saved formatted pptx: {saved_path}")
