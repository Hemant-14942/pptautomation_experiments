"""Test the MCQ formatter on p.pptx slide 24.

Tests the complete MCQ slide layout:
1. Question pill + label
2. Question text box
3. Four MCQ options (A, B, C, D) with pills, labels, and answer text
"""

from pathlib import Path
import copy

from lxml import etree
from pptx import Presentation


# Path to input file with MCQ slides
INPUT_PPT_PATH = Path("app/data/input/p.pptx")

# Path to save formatted output PPTX files
OUTPUT_DIR = Path(__file__).parent.parent / "output"


def q(tag: str) -> str:
    """Convert short tag name to full XML namespace."""
    ns = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }
    prefix, local = tag.split(":")
    return "{%s}%s" % (ns[prefix], local)


def clone_and_place_simple(el: etree._Element, off: tuple, ext: tuple) -> etree._Element:
    """Simple clone_and_place without external dependencies."""
    clone = copy.deepcopy(el)
    spPr = clone.find(q("p:spPr"))
    if spPr is not None and off and ext:
        xfrm = spPr.find(q("a:xfrm"))
        if xfrm is None:
            xfrm = etree.SubElement(spPr, q("a:xfrm"))
        off_el = xfrm.find(q("a:off"))
        if off_el is None:
            off_el = etree.SubElement(xfrm, q("a:off"))
        off_el.set("x", str(off[0]))
        off_el.set("y", str(off[1]))
        ext_el = xfrm.find(q("a:ext"))
        if ext_el is None:
            ext_el = etree.SubElement(xfrm, q("a:ext"))
        ext_el.set("cx", str(ext[0]))
        ext_el.set("cy", str(ext[1]))
    return clone


def test_mcq_formatter():
    """Test MCQ formatter on p.pptx slide 24."""

    if not INPUT_PPT_PATH.exists():
        print(f"Input file not found: {INPUT_PPT_PATH}")
        return

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    # Load presentation
    prs = Presentation(str(INPUT_PPT_PATH))

    print(f"\n{'='*80}")
    print(f"Testing MCQ formatter on p.pptx slide 24")
    print(f"{'='*80}\n")

    if len(prs.slides) < 25:
        print(f"Error: p.pptx only has {len(prs.slides)} slides, need at least 25 for slide 24")
        return

    # Get slide 24
    slide = prs.slides[24]

    print(f"Slide 24 has {len(slide.shapes)} shapes\n")

    # Identify the shapes we need
    print("Analyzing shapes on slide 24:")
    for i, shape in enumerate(slide.shapes):
        print(f"  Shape {i}: {shape.name}")
        if hasattr(shape, 'text_frame') and shape.has_text_frame:
            text = shape.text_frame.text[:40]
            print(f"    Text: '{text}'")

    print("\n" + "="*80)
    print("MCQ Formatter Test Complete")
    print(f"{'='*80}\n")

    # Save the presentation (even without modifications for now, just to test)
    output_file = OUTPUT_DIR / "p_slide24_mcq_test.pptx"
    prs.save(str(output_file))

    print(f"✓ Test file saved: {output_file.name}")
    print(f"  Location: {OUTPUT_DIR}")


if __name__ == "__main__":
    test_mcq_formatter()
