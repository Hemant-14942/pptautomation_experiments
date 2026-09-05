"""Diagnostic check: does a given (pptx_path, slide_index) qualify as
title_body_single_image, and if so, does the formatter place its body text
and image correctly into the fixed layout boxes?

Run directly for a manual check on one slide:
    python -m app.core.fit_layout.layout_formatter.title_body_single_image.tests.test_formatter <pptx_path> <slide_index>
"""

import os
import sys

from pptx import Presentation
from pptx.util import Emu

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "output")

from app.core.slide_text_analyzer import extract_text_from_slide
from app.core.style_parser import parse_input_slide_signature
from app.core.fit_layout.layout_formatter.slide_type_detector import (
    IMAGE_Y_THRESHOLD_EMU,
    SLIDE_TYPE_TITLE_BODY_SINGLE_IMAGE,
    detect_slide_type,
)
from app.core.fit_layout.layout_formatter.title_body_single_image.constants import (
    BODY_TEXTBOX,
    IMAGE_TEXTBOX,
)
from app.core.fit_layout.layout_formatter.title_body_single_image.formatter import (
    format_title_body_single_image,
)

XFRM_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def _get_off_ext(el):
    xfrm = el.find(".//p:spPr/a:xfrm", XFRM_NS)
    off = xfrm.find("a:off", XFRM_NS)
    ext = xfrm.find("a:ext", XFRM_NS)
    return (int(off.get("x")), int(off.get("y"))), (int(ext.get("cx")), int(ext.get("cy")))


def _find_body_and_picture(slide, prs):
    """Body text shape (>30 chars, excluding the heading's own shape) and
    the single real content picture (below the header area)."""
    heading = extract_text_from_slide(slide, 0, prs).get("heading")
    heading_pos = heading.get("position_top_left") if heading else None
    body_el = None
    pic_el = None
    pic_size = None
    for sp in slide.shapes:
        stype = sp.shape_type
        if stype is not None and getattr(stype, "name", "") == "PICTURE" and sp.top and sp.top > IMAGE_Y_THRESHOLD_EMU:
            pic_el = sp._element
            pic_size = sp.image.size
        elif sp.has_text_frame and (sp.left, sp.top) != heading_pos and len(sp.text_frame.text.strip()) > 30:
            body_el = sp._element
    return body_el, pic_el, pic_size


def verify_slide(pptx_path: str, slide_index: int) -> dict:
    """Check one slide end-to-end: detection + formatter placement.

    Returns a report dict — "pass": bool plus details for inspection.
    Does not raise on a non-matching slide; it reports why.
    """
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_index]
    sig = parse_input_slide_signature(pptx_path)[slide_index]

    detected_type = detect_slide_type(sig, slide, prs)
    if detected_type != SLIDE_TYPE_TITLE_BODY_SINGLE_IMAGE:
        return {
            "pass": False,
            "reason": f"slide type detected as {detected_type!r}, not title_body_single_image",
            "detected_type": detected_type,
        }

    body_el, pic_el, pic_size = _find_body_and_picture(slide, prs)

    if body_el is None or pic_el is None:
        return {
            "pass": False,
            "reason": "detector matched but could not re-locate body/picture shape",
            "detected_type": detected_type,
        }

    body_clone, pic_clone = format_title_body_single_image(
        body_el, pic_el, pic_size[0], pic_size[1]
    )

    body_off, body_ext = _get_off_ext(body_clone)
    pic_off, pic_ext = _get_off_ext(pic_clone)

    body_ok = (body_off, body_ext) == (
        (BODY_TEXTBOX["x"], BODY_TEXTBOX["y"]),
        (BODY_TEXTBOX["width"], BODY_TEXTBOX["height"]),
    )

    expected_aspect = pic_size[0] / pic_size[1]
    actual_aspect = pic_ext[0] / pic_ext[1]
    aspect_ok = abs(actual_aspect - expected_aspect) < 0.01

    within_box = (
        pic_off[0] >= IMAGE_TEXTBOX["x"]
        and pic_off[1] >= IMAGE_TEXTBOX["y"]
        and pic_off[0] + pic_ext[0] <= IMAGE_TEXTBOX["x"] + IMAGE_TEXTBOX["width"]
        and pic_off[1] + pic_ext[1] <= IMAGE_TEXTBOX["y"] + IMAGE_TEXTBOX["height"]
    )

    return {
        "pass": body_ok and aspect_ok and within_box,
        "detected_type": detected_type,
        "body": {"off": body_off, "ext": body_ext, "matches_fixed_box": body_ok},
        "image": {
            "off": pic_off,
            "ext": pic_ext,
            "aspect_ratio_preserved": aspect_ok,
            "within_image_box": within_box,
        },
    }


def save_formatted_slide(pptx_path: str, slide_index: int, output_dir: str = OUTPUT_DIR) -> str | None:
    """Apply the formatter to one slide in-place and save the whole
    presentation as a new pptx in output_dir, so the result can be opened
    and visually checked. Returns the saved path, or None if the slide
    doesn't match title_body_single_image."""
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_index]
    sig = parse_input_slide_signature(pptx_path)[slide_index]

    if detect_slide_type(sig, slide, prs) != SLIDE_TYPE_TITLE_BODY_SINGLE_IMAGE:
        return None

    body_el, pic_el, pic_size = _find_body_and_picture(slide, prs)

    if body_el is None or pic_el is None:
        return None

    body_clone, pic_clone = format_title_body_single_image(
        body_el, pic_el, pic_size[0], pic_size[1]
    )

    spTree = slide.shapes._spTree
    spTree.remove(body_el)
    spTree.remove(pic_el)
    spTree.append(body_clone)
    spTree.append(pic_clone)

    os.makedirs(output_dir, exist_ok=True)
    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    out_path = os.path.join(output_dir, f"{base_name}_slide{slide_index + 1}_formatted.pptx")
    prs.save(out_path)
    return out_path


def _print_report(pptx_path: str, slide_index: int, report: dict) -> None:
    print(f"\n{pptx_path} | slide index {slide_index} (PPT slide {slide_index + 1})")
    print(f"  detected_type: {report.get('detected_type')}")
    if not report["pass"] and "reason" in report:
        print(f"  RESULT: FAIL — {report['reason']}")
        return

    b = report["body"]
    i = report["image"]
    b_off_in = (Emu(b["off"][0]).inches, Emu(b["off"][1]).inches)
    b_ext_in = (Emu(b["ext"][0]).inches, Emu(b["ext"][1]).inches)
    i_off_in = (Emu(i["off"][0]).inches, Emu(i["off"][1]).inches)
    i_ext_in = (Emu(i["ext"][0]).inches, Emu(i["ext"][1]).inches)

    print(f"  body:  off={b_off_in[0]:.2f},{b_off_in[1]:.2f}in  ext={b_ext_in[0]:.2f}x{b_ext_in[1]:.2f}in  fixed_box_match={b['matches_fixed_box']}")
    print(f"  image: off={i_off_in[0]:.2f},{i_off_in[1]:.2f}in  ext={i_ext_in[0]:.2f}x{i_ext_in[1]:.2f}in  aspect_ok={i['aspect_ratio_preserved']}  within_box={i['within_image_box']}")
    print(f"  RESULT: {'PASS' if report['pass'] else 'FAIL'}")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python -m ...test_formatter <pptx_path> <slide_index>")
        sys.exit(1)
    path, idx = sys.argv[1], int(sys.argv[2])
    result = verify_slide(path, idx)
    _print_report(path, idx, result)
    if result["pass"]:
        saved_path = save_formatted_slide(path, idx)
        print(f"  saved formatted pptx: {saved_path}")
