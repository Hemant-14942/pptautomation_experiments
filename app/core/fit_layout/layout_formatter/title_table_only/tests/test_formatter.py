"""Diagnostic check: does a given (pptx_path, slide_index) qualify as
title_table_only, and if so, does the formatter place its table correctly
(centered if it fits, top-anchored if it doesn't) without resizing it?

Run directly for a manual check on one slide:
    python -m app.core.fit_layout.layout_formatter.title_table_only.tests.test_formatter <pptx_path> <slide_index>
"""

import os
import sys

from pptx import Presentation
from pptx.util import Emu

from app.core.style_parser import parse_input_slide_signature
from app.core.fit_layout.layout_formatter.slide_type_detector import (
    SLIDE_TYPE_TITLE_TABLE_ONLY,
    detect_slide_type,
)
from app.core.fit_layout.layout_formatter.title_table_only.constants import TABLE_BOX
from app.core.fit_layout.layout_formatter.title_table_only.formatter import (
    calculate_table_x,
    calculate_table_y,
    format_table_only,
)
from app.core.fit_layout.layout_formatter.title_table_only.table_utils import (
    get_table_total_height,
    get_table_width,
)

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "output")

XFRM_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def _get_off_ext(el):
    xfrm = el.find("p:xfrm", XFRM_NS)
    off = xfrm.find("a:off", XFRM_NS)
    ext = xfrm.find("a:ext", XFRM_NS)
    return (int(off.get("x")), int(off.get("y"))), (int(ext.get("cx")), int(ext.get("cy")))


def _find_table_el(slide):
    for sp in slide.shapes:
        if sp.has_table:
            return sp._element
    return None


def verify_slide(pptx_path: str, slide_index: int) -> dict:
    """Check one slide end-to-end: detection + formatter placement.

    Returns a report dict — "pass": bool plus details for inspection.
    Does not raise on a non-matching slide; it reports why.
    """
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_index]
    sig = parse_input_slide_signature(pptx_path)[slide_index]

    detected_type = detect_slide_type(sig, slide, prs)
    if detected_type != SLIDE_TYPE_TITLE_TABLE_ONLY:
        return {
            "pass": False,
            "reason": f"slide type detected as {detected_type!r}, not title_table_only",
            "detected_type": detected_type,
        }

    table_el = _find_table_el(slide)
    if table_el is None:
        return {
            "pass": False,
            "reason": "detector matched but could not locate table shape",
            "detected_type": detected_type,
        }

    orig_off, orig_ext = _get_off_ext(table_el)
    table_height = get_table_total_height(table_el)
    table_width = get_table_width(table_el)
    expected_y = calculate_table_y(table_height, TABLE_BOX)
    expected_x = calculate_table_x(table_width, TABLE_BOX)

    clone = format_table_only(table_el)
    new_off, new_ext = _get_off_ext(clone)

    position_ok = new_off == (expected_x, expected_y)
    size_unchanged = new_ext == orig_ext
    fits_in_box = table_height <= TABLE_BOX["height"]

    return {
        "pass": position_ok and size_unchanged,
        "detected_type": detected_type,
        "table_height": table_height,
        "fits_in_box": fits_in_box,
        "placement": "centered" if fits_in_box else "top-anchored",
        "off": new_off,
        "ext": new_ext,
        "orig_ext": orig_ext,
        "position_ok": position_ok,
        "size_unchanged": size_unchanged,
    }


def save_formatted_slide(pptx_path: str, slide_index: int, output_dir: str = OUTPUT_DIR) -> str | None:
    """Apply the formatter to one slide in-place and save the whole
    presentation as a new pptx in output_dir. Returns the saved path, or
    None if the slide doesn't match title_table_only."""
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_index]
    sig = parse_input_slide_signature(pptx_path)[slide_index]

    if detect_slide_type(sig, slide, prs) != SLIDE_TYPE_TITLE_TABLE_ONLY:
        return None

    table_el = _find_table_el(slide)
    if table_el is None:
        return None

    clone = format_table_only(table_el)
    spTree = slide.shapes._spTree
    spTree.remove(table_el)
    spTree.append(clone)

    os.makedirs(output_dir, exist_ok=True)
    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    out_path = os.path.join(output_dir, f"{base_name}_slide{slide_index + 1}_formatted.pptx")
    prs.save(out_path)
    return out_path


def _print_report(pptx_path: str, slide_index: int, report: dict) -> None:
    print(f"\n{pptx_path} | slide index {slide_index} (PPT slide {slide_index + 1})")
    print(f"  detected_type: {report.get('detected_type')}")
    if not report["pass"] and "reason" in report:
        print(f"  RESULT: FAIL — {report['reason']}")
        return

    h_in = Emu(report["table_height"]).inches
    off_in = (Emu(report["off"][0]).inches, Emu(report["off"][1]).inches)
    ext_in = (Emu(report["ext"][0]).inches, Emu(report["ext"][1]).inches)
    orig_ext_in = (Emu(report["orig_ext"][0]).inches, Emu(report["orig_ext"][1]).inches)

    print(f"  table height: {h_in:.2f}in  fits_in_box={report['fits_in_box']}  placement={report['placement']}")
    print(f"  off={off_in[0]:.2f},{off_in[1]:.2f}in  ext={ext_in[0]:.2f}x{ext_in[1]:.2f}in  (orig ext={orig_ext_in[0]:.2f}x{orig_ext_in[1]:.2f}in)")
    print(f"  position_ok={report['position_ok']}  size_unchanged={report['size_unchanged']}")
    print(f"  RESULT: {'PASS' if report['pass'] else 'FAIL'}")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python -m ...test_formatter <pptx_path> <slide_index>")
        sys.exit(1)
    path, idx = sys.argv[1], int(sys.argv[2])
    result = verify_slide(path, idx)
    _print_report(path, idx, result)
    if result["pass"]:
        saved_path = save_formatted_slide(path, idx)
        print(f"  saved formatted pptx: {saved_path}")
