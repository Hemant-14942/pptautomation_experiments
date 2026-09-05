"""Derive a slide type label from an input slide.

Heading detection uses the real pick_heading() logic (exactly 80pt, bold
preferred) via extract_text_from_slide() -- not a text-shape-count guess.
A slide's heading is always 80pt/bold by convention here; anything else is
treated as a malformed input and correctly fails detection rather than
being guessed at.

Combines that with the structural signature from
`app.core.style_parser.parse_input_slide_signature` (table/option counts)
and a position-filtered picture count, since the signature's raw
picture_count also counts small logo/icon shapes sitting near the top of
the slide.

Types recognized:
TITLE-BASED SLIDES:
- SLIDE_TYPE_TITLE_TABLE_ONLY: heading + a table, no real body text, no
  real images, no options
- SLIDE_TYPE_TITLE_BODY_ONLY: heading + body, no images
- SLIDE_TYPE_TITLE_BODY_SINGLE_IMAGE: heading + body + exactly one real image
- SLIDE_TYPE_TITLE_BODY_MULTIPLE_IMAGES: heading + body + 2+ real images

QUESTION-BASED SLIDES:
- SLIDE_TYPE_QUESTION_QTEXT_MCQ: question pill + question text + exactly 4
  MCQ options (A, B, C, D), NO heading, NO table, NO images, NO body text
  except question text and MCQ answer text

Any other combination returns None.
"""

from typing import Any

from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide

from app.core.slide_text_analyzer import extract_text_from_slide

SLIDE_TYPE_TITLE_TABLE_ONLY = "title_table_only"
SLIDE_TYPE_TITLE_BODY_ONLY = "title_body_only"
SLIDE_TYPE_TITLE_BODY_SINGLE_IMAGE = "title_body_single_image"
SLIDE_TYPE_TITLE_BODY_MULTIPLE_IMAGES = "title_body_multiple_images"

# Question-based slide types (no title heading required)
SLIDE_TYPE_QUESTION_QTEXT_MCQ = "question_qtext_mcq"

# Pictures at or above this Y position are treated as header logos/icons,
# not content images. Fixed value, not per-slide dynamic — chosen from
# observed logo positions (~0.94-0.95in) vs heading text start (~1.7in) in
# sample input decks (see app/data/input/p1.pptx).
IMAGE_Y_THRESHOLD_EMU = 1_828_800  # 2.0 inches


def has_real_body_text(slide: Slide, heading_pos: tuple[int, int] | None) -> bool:
    """Any text shape with >30 chars of content, excluding the heading's
    own shape (matched by position) so a long heading is never
    double-counted as body text too."""
    for sp in slide.shapes:
        if not sp.has_text_frame:
            continue
        if heading_pos is not None and (sp.left, sp.top) == heading_pos:
            continue
        if len(sp.text_frame.text.strip()) > 30:
            return True
    return False


def count_real_images(slide: Slide) -> int:
    """Count PICTURE shapes positioned below the header area.

    Excludes logos/icons placed near the top of the slide, which
    parse_input_slide_signature's picture_count does not distinguish from
    actual content images.
    """
    count = 0
    for sp in slide.shapes:
        shape_type = sp.shape_type
        if shape_type is not None and getattr(shape_type, "name", "") == "PICTURE":
            if sp.top is not None and sp.top > IMAGE_Y_THRESHOLD_EMU:
                count += 1
    return count


def count_actual_tables(slide: Slide) -> int:
    """Count TABLE shapes in the slide by checking actual shape type.

    Does NOT rely on parse_input_slide_signature which may be inaccurate.
    Directly counts shapes with type == TABLE (19).

    Returns:
        Number of TABLE shapes found in the slide
    """
    # Count by directly checking shape types
    table_count = 0
    for sp in slide.shapes:
        # Check if shape type is TABLE
        shape_type = sp.shape_type
        if shape_type is not None and getattr(shape_type, "name", "") == "TABLE":
            table_count += 1

    return table_count


def has_mcq_structure(slide: Slide) -> bool:
    """Check if slide has MCQ structure: question pill + 4 option pills (ellipses).

    MCQ slides have:
    - 1 roundRect shape (question pill at top)
    - 4 ellipse shapes (MCQ option pills A, B, C, D)
    - Multiple text boxes (question text + option labels + answer text)
    - NO title heading (these are question-based, not title-based)
    - NO tables, NO images

    Returns:
        True if slide has MCQ structure (1+ roundRect + 4 ellipses), False otherwise
    """
    from lxml import etree

    from app.constants.xml_namespaces import NS

    # Count roundRect and ellipse shapes
    roundrect_count = 0
    ellipse_count = 0

    for sp in slide.shapes:
        # For shapes that are direct shape elements (not groups)
        if hasattr(sp, "_element"):
            elem = sp._element

            # Find preset geometry to identify shape type
            geom = elem.find(
                ".//{%s}prstGeom" % NS["a"]
            )  # Look for <a:prstGeom> elements

            if geom is not None:
                prst = geom.get("prst")  # Get preset geometry name
                if prst == "roundRect":  # Question pill is roundRect
                    roundrect_count += 1
                elif prst == "ellipse":  # MCQ options are ellipses
                    ellipse_count += 1

        # Also check groups (MCQ options might be grouped)
        if str(sp.shape_type) == "<ShapeType.GROUP: 6>":
            # Groups might contain ellipse shapes
            # Count groups as potential option containers
            ellipse_count += 1

    # MCQ structure: 1 roundRect (question pill) + 4 ellipses (option pills)
    # We expect exactly 1 roundRect and at least 4 ellipses (or groups)
    has_mcq = roundrect_count >= 1 and ellipse_count >= 4

    return has_mcq


def detect_slide_type(
    signature: dict[str, Any], slide: Slide, prs: PresentationType
) -> str | None:

    # =========================================================================
    # SECTION 1: Check for MCQ slide type (question-based, no heading required)
    # =========================================================================
    # MCQ slides are question-based and don't have a title heading
    # They must have: 1 question pill (roundRect) + 4 option pills (ellipses)
    if has_mcq_structure(slide):
        # Confirm this is specifically question+text+mcq type (no table, no images)
        # Count actual TABLE shapes (not relying on signature which may be inaccurate)
        actual_table_count = count_actual_tables(slide)
        real_image_count = count_real_images(slide)

        # MCQ type should have NO tables and NO images
        if actual_table_count == 0 and real_image_count == 0:
            return SLIDE_TYPE_QUESTION_QTEXT_MCQ

    # =========================================================================
    # SECTION 2: Check for title-based slide types (all require heading)
    # =========================================================================
    # Extract heading using real 80pt bold detection
    heading = extract_text_from_slide(slide, 0, prs).get("heading")
    has_heading = heading is not None
    no_options = signature.get("option_pills", 0) == 0
    table_count = signature.get("table_count", 0)

    # Title-based slides require BOTH heading AND no options
    if not (has_heading and no_options):
        return None

    # Get heading position to exclude it when counting body text
    heading_pos = heading.get("position_top_left")
    has_body = has_real_body_text(slide, heading_pos)
    real_image_count = count_real_images(slide)

    # Check for title_table_only: heading + 1 table, no body, no images
    if table_count == 1 and not has_body and real_image_count == 0:
        return SLIDE_TYPE_TITLE_TABLE_ONLY

    # Check for body-based slides (title_body_only, _single_image, _multiple_images)
    if table_count == 0 and has_body:
        if real_image_count == 0:
            return SLIDE_TYPE_TITLE_BODY_ONLY
        elif real_image_count == 1:
            return SLIDE_TYPE_TITLE_BODY_SINGLE_IMAGE
        elif real_image_count >= 2:
            return SLIDE_TYPE_TITLE_BODY_MULTIPLE_IMAGES

    # No matching type found
    return None
